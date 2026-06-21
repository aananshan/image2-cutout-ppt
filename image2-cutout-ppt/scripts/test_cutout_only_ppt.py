from __future__ import annotations

import importlib.util
import json
import tempfile
import unittest
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


SCRIPT = Path(__file__).with_name("build_cutout_only_ppt.py")


def load_module():
    spec = importlib.util.spec_from_file_location("build_cutout_only_ppt", SCRIPT)
    module = importlib.util.module_from_spec(spec)
    assert spec.loader is not None
    spec.loader.exec_module(module)
    return module


class CutoutOnlyPptTests(unittest.TestCase):
    def test_builds_cutout_library_without_diagram_redrawing(self):
        mod = load_module()
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            asset_dir = root / "cutouts"
            asset_dir.mkdir()
            for name, color in [("cell", (80, 170, 220, 255)), ("drug", (45, 145, 75, 255))]:
                img = Image.new("RGBA", (180, 120), (0, 0, 0, 0))
                draw = ImageDraw.Draw(img)
                draw.ellipse((20, 16, 160, 104), fill=color, outline=(35, 70, 90, 255), width=4)
                img.save(asset_dir / f"{name}.png")
            manifest = asset_dir / "manifest.json"
            manifest.write_text(
                json.dumps(
                    [
                        {"id": "cell", "path": "cutouts/cell.png", "size_px": [180, 120]},
                        {"id": "drug", "path": "cutouts/drug.png", "size_px": [180, 120]},
                    ],
                    indent=2,
                ),
                encoding="utf-8",
            )
            out = mod.build_ppt(manifest, root / "exports" / "assets.pptx", project_root=root)

            self.assertTrue(out.exists())
            prs = Presentation(str(out))
            records = mod.load_records(manifest, project_root=root)
            pictures = [
                shape
                for slide in prs.slides
                for shape in slide.shapes
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            ]
            non_text_non_picture = [
                shape
                for slide in prs.slides
                for shape in slide.shapes
                if shape.shape_type != MSO_SHAPE_TYPE.PICTURE
                and not getattr(shape, "has_text_frame", False)
            ]

            self.assertGreaterEqual(len(pictures), len(records))
            self.assertEqual([], non_text_non_picture)
            self.assertGreaterEqual(len(prs.slides), 2)


if __name__ == "__main__":
    unittest.main()
