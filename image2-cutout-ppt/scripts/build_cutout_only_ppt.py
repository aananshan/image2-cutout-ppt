from __future__ import annotations

import json
import argparse
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx.dml.color import RGBColor
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DEFAULT_MANIFEST = ROOT / "cutouts" / "framed_sheet_v2" / "manifest.json"
DEFAULT_OUT = ROOT / "exports" / "image2_framed_cutouts_only_assets_v1.pptx"
DEFAULT_PREVIEW = ROOT / "analysis" / "image2_framed_cutouts_only_assets_v1_preview.png"


LABELS = {
    "lps": "LPS trigger",
    "ros": "ROS particles",
    "tlr4": "TLR4 receptor",
    "mitochondrion": "Mitochondrion",
    "nlrp3": "NLRP3 complex",
    "cytokine_particles": "Cytokine particles",
    "drug": "Drug capsule",
    "barrier_injury": "Barrier injury",
    "cytokine_storm": "Cytokine storm",
    "tissue_fibrosis": "Tissue fibrosis",
}


def load_records(manifest: Path = DEFAULT_MANIFEST, project_root: Path | None = None) -> list[dict]:
    manifest = Path(manifest)
    root = Path(project_root) if project_root is not None else manifest.parents[1]
    records = json.loads(manifest.read_text(encoding="utf-8"))
    for rec in records:
        rec["label"] = LABELS.get(rec["id"], rec["id"].replace("_", " ").title())
        raw_path = Path(rec["path"])
        rec["image_path"] = raw_path if raw_path.is_absolute() else root / raw_path
    return records


def image_size(path: Path) -> tuple[int, int]:
    with Image.open(path) as img:
        return img.size


def fit_box(path: Path, x: float, y: float, w: float, h: float) -> tuple[float, float, float, float]:
    iw, ih = image_size(path)
    scale = min(w / iw, h / ih)
    pw, ph = iw * scale, ih * scale
    return x + (w - pw) / 2, y + (h - ph) / 2, pw, ph


def add_text(slide, text: str, x: float, y: float, w: float, h: float, size: float, bold: bool = False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(24, 33, 52)
    return box


def place_asset(slide, rec: dict, x: float, y: float, w: float, h: float, label: bool = True):
    path = Path(rec["image_path"])
    px, py, pw, ph = fit_box(path, x, y, w, h - (0.26 if label else 0))
    slide.shapes.add_picture(str(path), Inches(px), Inches(py), Inches(pw), Inches(ph))
    if label:
        add_text(slide, rec["label"], x, y + h - 0.22, w, 0.18, 6.5)


def build_ppt(manifest: Path = DEFAULT_MANIFEST, out: Path = DEFAULT_OUT, project_root: Path | None = None) -> Path:
    records = load_records(manifest, project_root=project_root)
    out.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    add_text(slide, "Image2 Cutout Asset Library", 0.35, 0.18, 12.6, 0.34, 15, True)
    add_text(slide, "Each pattern is an independent transparent PNG object. Move, scale, recolor labels, and compose with native PPT arrows/text.", 0.9, 0.52, 11.5, 0.2, 7)

    cols = 5
    cell_w, cell_h = 2.38, 1.48
    start_x, start_y = 0.45, 0.95
    gap_x, gap_y = 0.22, 0.62
    for i, rec in enumerate(records):
        col = i % cols
        row = i // cols
        place_asset(
            slide,
            rec,
            start_x + col * (cell_w + gap_x),
            start_y + row * (cell_h + gap_y),
            cell_w,
            cell_h,
        )

    detail = prs.slides.add_slide(blank)
    add_text(detail, "Large Transparent Cutouts", 0.35, 0.18, 12.6, 0.34, 15, True)
    add_text(detail, "Same assets at larger size for checking edge completeness and copying into other slides.", 0.9, 0.52, 11.5, 0.2, 7)
    big_positions = [
        (0.55, 0.95, 2.35, 1.8),
        (3.15, 0.95, 2.55, 1.8),
        (5.95, 0.9, 2.75, 1.9),
        (9.0, 0.85, 3.15, 1.95),
        (0.65, 3.2, 2.8, 2.2),
        (3.9, 3.35, 2.5, 1.9),
        (6.8, 3.15, 2.75, 2.0),
        (9.9, 3.05, 2.75, 2.1),
    ]
    for rec, pos in zip(records, big_positions):
        place_asset(detail, rec, *pos)

    if len(records) > len(big_positions):
        extra = prs.slides.add_slide(blank)
        add_text(extra, "Additional Cutouts", 0.35, 0.18, 12.6, 0.34, 15, True)
        for i, rec in enumerate(records[len(big_positions) :]):
            place_asset(extra, rec, 0.8 + i * 3.0, 1.2, 2.65, 2.2)

    prs.save(out)
    return out


def build_preview(manifest: Path = DEFAULT_MANIFEST, out: Path = DEFAULT_PREVIEW, project_root: Path | None = None) -> Path:
    records = load_records(manifest, project_root=project_root)
    out.parent.mkdir(parents=True, exist_ok=True)
    canvas = Image.new("RGB", (1600, 900), "white")
    draw = ImageDraw.Draw(canvas)
    try:
        title_font = ImageFont.truetype("arial.ttf", 34)
        label_font = ImageFont.truetype("arial.ttf", 16)
    except OSError:
        title_font = ImageFont.load_default()
        label_font = ImageFont.load_default()
    draw.text((44, 28), "Image2 Cutout Asset Library", fill=(24, 33, 52), font=title_font)
    draw.text((44, 75), "Independent transparent PNG objects arranged for copying into PPT.", fill=(72, 82, 101), font=label_font)

    cols = 5
    tile_w, tile_h = 285, 230
    sx, sy = 55, 135
    gx, gy = 25, 38
    for i, rec in enumerate(records):
        x = sx + (i % cols) * (tile_w + gx)
        y = sy + (i // cols) * (tile_h + gy)
        cut = Image.open(rec["image_path"]).convert("RGBA")
        max_w, max_h = tile_w - 22, tile_h - 48
        scale = min(max_w / cut.width, max_h / cut.height)
        cut = cut.resize((max(1, round(cut.width * scale)), max(1, round(cut.height * scale))), Image.Resampling.LANCZOS)
        canvas.paste(cut, (x + (tile_w - cut.width) // 2, y + (max_h - cut.height) // 2), cut)
        tw = draw.textlength(rec["label"], font=label_font)
        draw.text((x + (tile_w - tw) / 2, y + tile_h - 28), rec["label"], fill=(40, 48, 66), font=label_font)

    canvas.save(out)
    return out


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build a cutout-only PPT asset library from a manifest of transparent PNGs.")
    parser.add_argument("--manifest", type=Path, default=DEFAULT_MANIFEST, help="Path to cutout manifest.json.")
    parser.add_argument("--project-root", type=Path, default=None, help="Root used to resolve relative image paths. Defaults to manifest grandparent.")
    parser.add_argument("--out", type=Path, default=DEFAULT_OUT, help="Output PPTX path.")
    parser.add_argument("--preview", type=Path, default=DEFAULT_PREVIEW, help="Output PNG preview path.")
    parser.add_argument("--no-preview", action="store_true", help="Skip PNG preview generation.")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    pptx = build_ppt(args.manifest, args.out, project_root=args.project_root)
    preview = None if args.no_preview else build_preview(args.manifest, args.preview, project_root=args.project_root)
    print(
        json.dumps(
            {"pptx": str(pptx), "preview": str(preview) if preview else None, "count": len(load_records(args.manifest, project_root=args.project_root))},
            ensure_ascii=False,
            indent=2,
        )
    )


if __name__ == "__main__":
    main()
